#!/usr/bin/env python3
import argparse

from pptx import Presentation


def _squash_ws(s: str) -> str:
    return " ".join(s.split())


def main() -> int:
    parser = argparse.ArgumentParser(description="Print per-slide text from a PPTX.")
    parser.add_argument("pptx", help="Path to .pptx")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            t = _squash_ws(shape.text_frame.text or "")
            if t:
                texts.append(t)

        header = f"Slide {i:02d}"
        if not texts:
            print(f"{header}: (no text)")
            continue

        print(f"{header}:")
        for t in texts:
            print(f"  - {t}")
        print()

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
