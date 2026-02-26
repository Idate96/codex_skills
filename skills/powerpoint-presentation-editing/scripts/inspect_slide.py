#!/usr/bin/env python3
import argparse

from pptx import Presentation


def _short(s: str, max_len: int = 120) -> str:
    s = " ".join((s or "").split())
    if len(s) <= max_len:
        return s
    return s[: max_len - 3] + "..."


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Inspect shapes on a slide (type, bounds, and text preview)."
    )
    parser.add_argument("pptx", help="Path to .pptx")
    parser.add_argument("--slide", type=int, required=True, help="1-based slide index")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    if args.slide < 1 or args.slide > len(prs.slides):
        raise SystemExit(
            f"ERROR: --slide must be in [1, {len(prs.slides)}], got {args.slide}"
        )

    slide = prs.slides[args.slide - 1]
    print(f"Slide {args.slide:02d}: {len(slide.shapes)} shapes")

    for idx, shape in enumerate(slide.shapes, start=1):
        shape_type = getattr(shape, "shape_type", None)
        name = getattr(shape, "name", "")

        def _emu(attr: str) -> str:
            v = getattr(shape, attr, None)
            return str(v) if v is not None else "?"

        bounds = f"L={_emu('left')} T={_emu('top')} W={_emu('width')} H={_emu('height')}"
        line = f"{idx:02d}. type={shape_type} name={name!r} {bounds}"

        if getattr(shape, "has_text_frame", False):
            text = _short(getattr(shape, "text", ""))
            if text:
                line += f" text={text!r}"

        # Pictures: show the packaged filename if available (useful for tracking media).
        img = getattr(shape, "image", None)
        if img is not None and getattr(img, "filename", None):
            line += f" image={img.filename!r}"

        print(line)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
